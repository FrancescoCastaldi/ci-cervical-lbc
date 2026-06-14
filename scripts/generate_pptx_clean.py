"""Generate a clean, minimal PowerPoint presentation for CI Cervical LBC project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Constants ──────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Inches(0.9)
MARGIN_TOP = Inches(0.6)
CONTENT_WIDTH = Inches(11.5)

FONT = "Calibri"
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_SUBTLE = RGBColor(0x88, 0x88, 0x88)
COLOR_ACCENT = RGBColor(0x22, 0x22, 0x22)

OUTPUT = os.path.join(os.path.dirname(os.path.dirname(__file__)), "slides", "presentazione_v2.pptx")


def add_slide(prs):
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    return txBox.text_frame


def set_paragraph(tf, text, size=18, bold=False, color=COLOR_TEXT, alignment=PP_ALIGN.LEFT, space_after=Pt(6), space_before=Pt(0)):
    """Write text into a text frame, clearing existing content."""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_run(p, text, size=18, bold=False, color=COLOR_TEXT, italic=False):
    """Add a run to an existing paragraph."""
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def add_paragraph(tf, text="", size=18, bold=False, color=COLOR_TEXT, alignment=PP_ALIGN.LEFT,
                  space_after=Pt(4), space_before=Pt(0), bullet=False):
    """Add a new paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    if text:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return p


def slide_title_bar(slide, title, subtitle=None):
    """Add a clean title bar at the top of the slide."""
    # Title
    tf = add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7))
    tf.word_wrap = True
    set_paragraph(tf, title, size=30, bold=True, space_after=Pt(2))

    # Thin line under title
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, MARGIN_TOP + Inches(0.65),
        CONTENT_WIDTH, Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line.line.fill.background()

    if subtitle:
        tf2 = add_textbox(slide, MARGIN_LEFT, MARGIN_TOP + Inches(0.75), CONTENT_WIDTH, Inches(0.4))
        tf2.word_wrap = True
        set_paragraph(tf2, subtitle, size=14, color=COLOR_SUBTLE, space_after=Pt(0))

    return MARGIN_TOP + Inches(1.3)  # return Y position after title


def add_body_text(slide, text, y, size=16, line_spacing=1.5):
    """Add a multi-line body text."""
    tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5))
    tf.word_wrap = True
    lines = text.strip().split('\n')
    first = True
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if first:
            set_paragraph(tf, line, size=size, space_after=Pt(2))
            first = False
        else:
            add_paragraph(tf, line, size=size, space_after=Pt(2))
    return tf


def add_bullet_list(slide, items, y, size=16):
    """Add a bulleted list."""
    tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5))
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            set_paragraph(tf, f"  \u2022  {item}", size=size, space_after=Pt(3))
            first = False
        else:
            add_paragraph(tf, f"  \u2022  {item}", size=size, space_after=Pt(3))
    return tf


def add_table(slide, headers, rows, y, col_widths=None):
    """Add a clean table."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          MARGIN_LEFT, y,
                                          CONTENT_WIDTH, Inches(0.35 * n_rows))
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Style header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
            p.font.name = FONT

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = FONT
                p.font.color.rgb = COLOR_TEXT

    return y + Inches(0.35 * n_rows) + Inches(0.2)


# ══════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ── SLIDE 1: Title ────────────────────────────────────────────────────
slide = add_slide(prs)

# Vertical center
tf = add_textbox(slide, MARGIN_LEFT, Inches(2.2), CONTENT_WIDTH, Inches(3))
tf.word_wrap = True
set_paragraph(tf, "Deblur & Denoise di immagini LBC Cervical", size=36, bold=True,
              alignment=PP_ALIGN.CENTER, space_after=Pt(8))
add_paragraph(tf, "Confronto tra metodi variazionale, deep learning e generativo",
              size=20, color=COLOR_SUBTLE, alignment=PP_ALIGN.CENTER, space_after=Pt(30))
add_paragraph(tf, "Francesco Castaldi, Paolo Fusco", size=16,
              alignment=PP_ALIGN.CENTER, space_after=Pt(2))
add_paragraph(tf, "Computational Imaging \u2014 Universit\u00e0 di Bologna",
              size=14, color=COLOR_SUBTLE, alignment=PP_ALIGN.CENTER, space_after=Pt(2))
add_paragraph(tf, "Prof. Picciolomini & Evangelista", size=14, color=COLOR_SUBTLE,
              alignment=PP_ALIGN.CENTER, space_after=Pt(0))


# ── SLIDE 2: Description ──────────────────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "Descrizione del Progetto",
                    "Problema inverso: y = (k \u2217 x) + n \u2014 recuperare l\u2019immagine pulita x dall\u2019osservazione degradata y")

sections = [
    ("Degradazione", [
        "Blur Gaussiano: \u03c3=2, kernel 9\u00d79",
        "AWGN: \u03c3\u2099 \u2208 {0.005, 0.01, 0.05, 0.1}",
        "Stessa pipeline per tutti (seed=42)"
    ]),
    ("Dataset", [
        "Mendeley LBC Cervical Cancer",
        "962 immagini, 4 classi diagnostiche",
        "256\u00d7256, split 70/15/15"
    ]),
    ("Metodi", [
        "TV \u2014 Variazionale",
        "UNet \u2014 Deep Learning",
        "DiffPIR \u2014 Generativo (Diffusion)"
    ])
]

x_positions = [MARGIN_LEFT, MARGIN_LEFT + Inches(4.0), MARGIN_LEFT + Inches(8.0)]
for (title, items), x in zip(sections, x_positions):
    tf = add_textbox(slide, x, y, Inches(3.3), Inches(3))
    tf.word_wrap = True
    set_paragraph(tf, title, size=18, bold=True, space_after=Pt(8))
    for item in items:
        add_paragraph(tf, f"  \u2022  {item}", size=14, space_after=Pt(3))

# Bottom note
tf = add_textbox(slide, MARGIN_LEFT, Inches(5.8), CONTENT_WIDTH, Inches(0.6))
tf.word_wrap = True
set_paragraph(tf, "Obiettivo: Confronto equo su stessi input \u2014 PSNR, SSIM, tempi ai vari livelli di rumore",
              size=16, bold=True, alignment=PP_ALIGN.CENTER)


# ── SLIDE 3: Methodology ──────────────────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "3 Metodi, 3 Famiglie")

methods = [
    ("1. Total Variation (TV) \u2014 Variazionale",
     "Modello: minimizza J(x) = \u00bd\u2016Hx \u2212 y\u2016\u00b2 + \u03bb \u00b7 TV(x)\n"
     "TV(x) = \u03a3 \u221a(|\u2207_h x|\u00b2 + |\u2207_v x|\u00b2) \u2014 penalizza oscillazioni\n"
     "Ottimizzazione: Adam, 150 iter, lr=0.001, \u03bb=0.005",
     "Pro: nessun training, interpretabile | Contro: \u03bb fisso, staircasing"),
    ("2. UNet \u2014 End-to-End Deep Learning",
     "Encoder-decoder con skip connections (4 livelli, 64\u2192512 canali)\n"
     "Loss: MSE + Adam, lr=1e-4, batch=16. Multi-noise augmentation.\n"
     "Training: 50 epoche CPU, 1.9M params, GroupNorm, noise cond.",
     "Pro: inferenza 26 ms/img | Contro: richiede training e GPU"),
    ("3. DiffPIR \u2014 Generativo (Diffusion)",
     "DDPM: inverte forward diffusion, rimuove rumore gradualmente\n"
     "LightUNet: 1.26M params addestrata su LBC\n"
     "DiffPIR: data-fidelity FFT + DDIM. 15 steps, t_start=50, \u03bb=10",
     "Pro: eccelle ad alto rumore | Contro: lento (2s/img), allucina")
]

for title, body, procon in methods:
    tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(1.5))
    tf.word_wrap = True
    set_paragraph(tf, title, size=18, bold=True, space_after=Pt(4))

    lines = body.split('\n')
    for line in lines:
        add_paragraph(tf, line, size=13, space_after=Pt(2))

    add_paragraph(tf, procon, size=12, color=COLOR_SUBTLE, space_after=Pt(10), space_before=Pt(4))
    y += Inches(1.7)


# ── SLIDE 4: Implementation ───────────────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "Implementation")

# Pipeline
tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.8))
tf.word_wrap = True
set_paragraph(tf, "Pipeline comune", size=18, bold=True, space_after=Pt(4))
p = add_paragraph(tf, "", size=14, space_after=Pt(0))
parts = [
    ("Dataset", False), (" \u2192 ", True), ("resize 256\u00d7256", False),
    (" \u2192 ", True), ("norm [-1, 1]", False), (" \u2192 ", True),
    ("degradazione (blur \u03c3=2 + AWGN)", False), (" \u2192 ", True),
    ("metodo (TV / UNet / DiffPIR)", False), (" \u2192 ", True),
    ("valutazione (PSNR, SSIM)", False)
]
for text, is_arrow in parts:
    add_run(p, text, size=13, bold=not is_arrow, color=COLOR_ACCENT if is_arrow else COLOR_TEXT)

y += Inches(0.8)

# Three columns for methods
impl_data = [
    ("TV \u2014 Algoritmo",
     "1. Inizializza x\u2080 = y\n"
     "2. Calcola gradiente J(x)\n"
     "3. Adam: lr=0.001, 150 iter\n"
     "4. \u03bb=0.005 (euristico)"),
    ("UNet \u2014 Architettura",
     "Features: [16, 32, 64, 128]\n"
     "DoubleConv: Conv+GroupNorm+ReLU\n"
     "Loss: L1, Adam lr=1e-4\n"
     "50 epoche, batch=16, CPU"),
    ("DiffPIR \u2014 Algoritmo",
     "1. x\u209c\u2081 ~ N(0, I)\n"
     "2. Denoising: DDIM(x\u209c, t)\n"
     "3. Data-fidelity via FFT\n"
     "4. 15 steps, t_start=50")
]

cols_x = [MARGIN_LEFT, MARGIN_LEFT + Inches(4.0), MARGIN_LEFT + Inches(8.0)]
for (title, body), x in zip(impl_data, cols_x):
    tf = add_textbox(slide, x, y, Inches(3.3), Inches(3))
    tf.word_wrap = True
    set_paragraph(tf, title, size=16, bold=True, space_after=Pt(6))
    lines = body.split('\n')
    for line in lines:
        add_paragraph(tf, line, size=12, space_after=Pt(3))


# ── SLIDE 5: DiffPIR Implementation ──────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "DiffPIR \u2014 Implementazione",
                    "DDPM + FFT Data-Fidelity per deblur & denoise")

# Left column: code
left_x = MARGIN_LEFT
tf = add_textbox(slide, left_x, y, Inches(5.8), Inches(5.5))
tf.word_wrap = True

# Set default font to Consolas for code look
set_paragraph(tf, "Data-fidelity nel dominio FFT", size=14, bold=True, space_after=Pt(4))
code1 = (
    "def _data_fidelity_fft(y, x0_pred, kernel_otf, rho):\n"
    "    FB = kernel_otf\n"
    "    FBC = torch.conj(FB)\n"
    "    F2B = torch.abs(FB) ** 2\n"
    "    y_f = torch.fft.fftn(y, dim=(-2, -1))\n"
    "    x0_f = torch.fft.fftn(x0_pred, dim=(-2, -1))\n"
    "    num = FBC * y_f + rho * x0_f\n"
    "    den = F2B + rho\n"
    "    x = torch.fft.ifftn(num/den, dim=(-2, -1)).real\n"
    "    return x"
)
for line in code1.split('\n'):
    add_paragraph(tf, line, size=10, color=COLOR_TEXT, space_after=Pt(0))

add_paragraph(tf, "", size=4, space_after=Pt(4))
add_paragraph(tf, "Gaussian Blur Kernel (PSF) \u2192 OTF", size=14, bold=True, space_after=Pt(4))
code2 = (
    "def _psf_to_otf(kernel, img_size):\n"
    "    pad_h = img_size[-2] - kernel.shape[-2]\n"
    "    pad_w = img_size[-1] - kernel.shape[-1]\n"
    "    padded = F.pad(kernel, (pad_w//2, ...))\n"
    "    shifted = torch.roll(padded, ...)\n"
    "    otf = torch.fft.fftn(shifted, dim=(-2, -1))\n"
    "    return otf"
)
for line in code2.split('\n'):
    add_paragraph(tf, line, size=10, color=COLOR_TEXT, space_after=Pt(0))

add_paragraph(tf, "", size=4, space_after=Pt(4))
add_paragraph(tf, "Sampling loop (DDIM + Data-Fidelity)", size=14, bold=True, space_after=Pt(4))
code3 = (
    "# Subsampled timesteps: t_start \u2192 0\n"
    "for t_i in seq:\n"
    "    # 1. Denoising: predicted noise\n"
    "    pred_noise = model(x, t_tensor)\n"
    "    # 2. Estimate x0 from x_t\n"
    "    x0 = (x - sqrt_1m_ac[t_i] * pred) / sqrt_ac[t_i]\n"
    "    x0 = x0.clamp(-1, 1)\n"
    "    # 3. Data-fidelity (FFT correction)\n"
    "    x0 = _data_fidelity_fft(y, x0, otf, rho_t)\n"
    "    # 4. DDIM step: predict x_{t-1}\n"
    "    eps = (x - sqrt_ac[t_i] * x0) / sqrt_1m_ac[t_i]\n"
    "    x = sqrt_ac[t_{i-1}] * x0 + sqrt_1m_ac[t_{i-1}] * eps"
)
for line in code3.split('\n'):
    add_paragraph(tf, line, size=10, color=COLOR_TEXT, space_after=Pt(0))


# Right column: explanation + architecture
right_x = MARGIN_LEFT + Inches(6.5)
tf2 = add_textbox(slide, right_x, y, Inches(5), Inches(5.5))
tf2.word_wrap = True

set_paragraph(tf2, "Architettura LightUNet", size=16, bold=True, space_after=Pt(6))
arch_items = [
    "Encoder: 3 livelli (32, 64, 128 canali)",
    "Mid: 2 ResBlock (128 canali)",
    "Decoder: 3 livelli (skip connections)",
    "Sinusoidal time embedding",
    "1.26M parametri totali",
]
for item in arch_items:
    add_paragraph(tf2, f"  \u2022  {item}", size=13, space_after=Pt(2))

add_paragraph(tf2, "", size=6, space_after=Pt(6))
add_paragraph(tf2, "Parametri DiffPIR", size=16, bold=True, space_after=Pt(6))
params = [
    ("DDIM steps:", "15 (sub-campionati da t_start=50)"),
    ("Lambda (\u03bb):", "10.0 (peso data-fidelity)"),
    ("Zeta (\u03b6):", "0 (deterministico)"),
    ("t_start:", "50 (stabilit\u00e0 numerica)"),
    ("Noise levels:", "0.005 / 0.01 / 0.05 / 0.1"),
]
for label, val in params:
    add_paragraph(tf2, f"{label}  {val}", size=13, space_after=Pt(3))

add_paragraph(tf2, "", size=6, space_after=Pt(6))
add_paragraph(tf2, "Pesi del Modello", size=16, bold=True, space_after=Pt(6))
add_paragraph(tf2, "DDPM addestrato su 100 immagini LBC", size=13, space_after=Pt(2))
add_paragraph(tf2, "30 epoche, MSE loss, Adam lr=1e-4", size=13, space_after=Pt(2))
add_paragraph(tf2, "Pesi: ~5 MB (ddpm_lbc.pt)", size=13, color=COLOR_SUBTLE, space_after=Pt(2))


# ── SLIDE 7: Experimental Setup ──────────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "Impostazione Sperimentale",
                    "Deblur & Denoise \u2014 Immagini LBC Cervicali")

# Left column: problem + blur + noise
left_x = MARGIN_LEFT
tf = add_textbox(slide, left_x, y, Inches(5.5), Inches(4.5))
tf.word_wrap = True
set_paragraph(tf, "Problema inverso: y = (k \u2217 x) + n", size=18, bold=True, space_after=Pt(8))
add_paragraph(tf, "immagine degradata = blur(originale) + rumore", size=14, color=COLOR_SUBTLE, space_after=Pt(14))

add_paragraph(tf, "Kernel di Blur (PSF Gaussiana)", size=16, bold=True, space_after=Pt(4))
for line in ["\u03c3 = 2.0", "Dimensione: 9 \u00d7 9", "Normalizzazione: somma = 1"]:
    add_paragraph(tf, f"  \u2022  {line}", size=13, space_after=Pt(2))

add_paragraph(tf, "", size=8, space_after=Pt(4))
add_paragraph(tf, "AWGN \u2014 4 Livelli di Rumore", size=16, bold=True, space_after=Pt(4))
for line in ["0.005 \u2014 Molto basso", "0.01  \u2014 Basso", "0.05  \u2014 Medio", "0.1   \u2014 Alto"]:
    add_paragraph(tf, f"  \u2022  {line}", size=13, space_after=Pt(2))

add_paragraph(tf, "Seed=42 per tutti i metodi", size=12, color=COLOR_SUBTLE, space_after=Pt(0), space_before=Pt(8))

# Right column: dataset
right_x = MARGIN_LEFT + Inches(6.5)
tf = add_textbox(slide, right_x, y, Inches(5), Inches(4.5))
tf.word_wrap = True
set_paragraph(tf, "Dataset", size=18, bold=True, space_after=Pt(8))

dataset_info = [
    ("Fonte:", "Mendeley LBC Cervical"),
    ("Immagini:", "962 RGB"),
    ("Classi:", "NILM, HSIL, LSIL, SCC"),
    ("Dimensione:", "256\u00d7256, norm [-1, 1]"),
    ("Split:", "70/15/15"),
]
for label, value in dataset_info:
    add_paragraph(tf, f"{label}  {value}", size=14, space_after=Pt(4))

add_paragraph(tf, "", size=6, space_after=Pt(4))
add_paragraph(tf, "Distribuzione:", size=14, bold=True, space_after=Pt(4))
for cls, count, pct in [("NILM", 359, "37.3%"), ("HSIL", 247, "25.7%"),
                          ("LSIL", 231, "24.0%"), ("SCC", 125, "13.0%")]:
    add_paragraph(tf, f"  \u2022  {cls:6s}  {count:3d}  ({pct})", size=13, space_after=Pt(2))


# ── SLIDE 8: Results ─────────────────────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "Risultati Numerici",
                    "PSNR \u00b7 SSIM \u00b7 Tempo di inferenza \u2014 Test Set (145 immagini)")

# TV Table
tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.4))
set_paragraph(tf, "Total Variation (TV) \u2014 Baseline Variazionale", size=16, bold=True, space_after=Pt(4))
y += Inches(0.4)

headers = ["\u03c3\u2099", "PSNR", "SSIM"]
tv_rows = [["0.005", "32.09 dB", "0.911"], ["0.01", "32.04 dB", "0.909"],
           ["0.05", "30.42 dB", "0.837"], ["0.1", "26.54 dB", "0.586"]]
y = add_table(slide, headers, tv_rows, y, col_widths=[Inches(2), Inches(2), Inches(2)])

# UNet Table
tf = add_textbox(slide, MARGIN_LEFT, y + Inches(0.1), CONTENT_WIDTH, Inches(0.4))
set_paragraph(tf, "UNet \u2014 Deep Learning End-to-End", size=16, bold=True, space_after=Pt(4))
y += Inches(0.5)

headers = ["\u03c3\u2099", "PSNR", "SSIM", "Tempo"]
unet_rows = [["0.005", "29.89 dB", "0.894", "0.035 s"],
             ["0.01", "29.89 dB", "0.894", "0.034 s"],
             ["0.05", "29.63 dB", "0.875", "0.034 s"],
             ["0.1", "28.93 dB", "0.830", "0.036 s"]]
y = add_table(slide, headers, unet_rows, y, col_widths=[Inches(1.5), Inches(2), Inches(1.5), Inches(1.5)])

# DiffPIR Table
tf = add_textbox(slide, MARGIN_LEFT, y + Inches(0.1), CONTENT_WIDTH, Inches(0.4))
set_paragraph(tf, "DiffPIR \u2014 Generative (Diffusion)", size=16, bold=True, space_after=Pt(4))
y += Inches(0.5)

headers = ["\u03c3\u2099", "PSNR", "SSIM", "Tempo"]
diffpir_rows = [["0.005", "16.67 dB", "0.235", "3.27 s"],
                ["0.01", "17.32 dB", "0.270", "3.00 s"],
                ["0.05", "22.49 dB", "0.512", "2.85 s"],
                ["0.1", "24.68 dB", "0.663", "2.89 s"]]
y = add_table(slide, headers, diffpir_rows, y, col_widths=[Inches(1.5), Inches(2), Inches(1.5), Inches(1.5)])

# Observations
y += Inches(0.2)
tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(1.5))
tf.word_wrap = True
set_paragraph(tf, "Osservazioni", size=14, bold=True, space_after=Pt(4))
for obs in [
    "TV: eccelle a basso rumore (32 dB), degrada ad alto rumore (26.5 dB) \u2014 limite \u03bb fisso",
    "UNet: pi\u00f9 robusto, PSNR stabile 28.9\u201329.9 dB, inferenza ~0.035 s",
    "DiffPIR: migliora col rumore (+8 dB), richiede ~3 s/img",
    "Trade-off: UNet vince in velocit\u00e0/consistenza, DiffPIR in flessibilit\u00e0, TV in semplicit\u00e0",
]:
    add_paragraph(tf, f"  \u2022  {obs}", size=12, space_after=Pt(2))


# ── SLIDE 9: Conclusions ─────────────────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "Conclusioni",
                    "Riassunto, risultati e possibili sviluppi futuri")

tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(1.2))
tf.word_wrap = True
set_paragraph(tf, "Riassunto", size=18, bold=True, space_after=Pt(4))
add_paragraph(tf,
    "Abbiamo affrontato il problema inverso di deblur & denoise su immagini LBC cervicali "
    "(Mendeley dataset, 962 immagini, 256\u00d7256) confrontando tre metodi di natura diversa: "
    "Total Variation (variazionale), UNet (deep learning end-to-end) e "
    "DiffPIR (generativo con modelli di diffusione). Tutti i metodi condividono la stessa "
    "pipeline di degradazione (blur Gaussiano \u03c3=2 + AWGN a 4 livelli) per un confronto equo.",
    size=14, space_after=Pt(2))

y += Inches(1.3)

# Key results cards
card_data = [
    ("32.09 dB", "TV a \u03c3=0.005", "Miglior PSNR a basso rumore"),
    ("29.9 dB", "UNet (stabile)", "Robusto, 0.035 s"),
    ("24.68 dB", "DiffPIR a \u03c3=0.1", "Migliora con rumore alto, ~2 s/img"),
]

card_width = Inches(3.6)
card_gap = Inches(0.3)
total_cards = card_width * 3 + card_gap * 2
start_x = MARGIN_LEFT + (CONTENT_WIDTH - total_cards) // 2

for i, (value, label, desc) in enumerate(card_data):
    x = start_x + i * (card_width + card_gap)
    tf = add_textbox(slide, x, y, card_width, Inches(1.0))
    tf.word_wrap = True
    set_paragraph(tf, value, size=26, bold=True, alignment=PP_ALIGN.CENTER, space_after=Pt(2))
    add_paragraph(tf, label, size=14, color=COLOR_SUBTLE, alignment=PP_ALIGN.CENTER, space_after=Pt(2))
    add_paragraph(tf, desc, size=12, color=COLOR_SUBTLE, alignment=PP_ALIGN.CENTER, space_after=Pt(0))

y += Inches(1.3)

tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(2.5))
tf.word_wrap = True
set_paragraph(tf, "Osservazioni", size=16, bold=True, space_after=Pt(4))
for obs in [
    "TV: migliore a basso rumore (32.09 dB), ma \u03bb fisso penalizza ad alto rumore (26.54 dB).",
    "UNet: pi\u00f9 robusto e veloce. PSNR 28.9\u201329.9 dB, inferenza 0.035 s.",
    "DiffPIR: migliora col rumore (16.7 \u2192 24.7 dB), ma richiede ~2 s/img.",
    "Trade-off: UNet per velocit\u00e0, TV per semplicit\u00e0, DiffPIR per flessibilit\u00e0.",
]:
    add_paragraph(tf, f"  \u2022  {obs}", size=13, space_after=Pt(2))

add_paragraph(tf, "", size=6, space_after=Pt(4))
add_paragraph(tf, "Possibili Sviluppi Futuri", size=16, bold=True, space_after=Pt(4))
for fut in [
    "Dati reali: test su degradazioni non simulate",
    "TV adattivo: regolarizzazione dipendente dal noise level",
    "UNet + DiffPIR: UNet come initial guess per ridurre step diffusione",
    "Altri dataset medici: validare su altre immagini citologiche",
    "Segmentazione: integrare deblur/denoise con segmentazione cellulare",
]:
    add_paragraph(tf, f"  \u2022  {fut}", size=13, space_after=Pt(2))


# ── SLIDE 10: Bibliography ───────────────────────────────────────────
slide = add_slide(prs)
y = slide_title_bar(slide, "Bibliografia e Riferimenti",
                    "Articoli scientifici, dataset e strumenti utilizzati nel progetto")

# Single textbox with all content — no overlapping, auto-flow
tf = add_textbox(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.5))
tf.word_wrap = True

entries = [
    # (section_or_ref, is_section, name, authors, title_venue)
    (True,  "Metodi di Restauro", "", ""),
    (False, "Total Variation (TV)",
     "Rudin, L., Osher, S., Fatemi, E.",
     "\u201cNonlinear total variation based noise removal algorithms.\u201d \u2014 Physica D: Nonlinear Phenomena, 1992."),
    (False, "U-Net",
     "Ronneberger, O., Fischer, P., Brox, T.",
     "\u201cU-Net: Convolutional Networks for Biomedical Image Segmentation.\u201d \u2014 MICCAI 2015."),
    (False, "DiffPIR",
     "Zhu, Y., Zhang, K., Liang, J., Cao, J., Wen, B., Timofte, R., Van Gool, L.",
     "\u201cDenoising Diffusion Models for Plug-and-Play Image Restoration.\u201d \u2014 CVPR 2023, arXiv:2305.08995"),
    (True,  "Modelli di Diffusione", "", ""),
    (False, "DDPM",
     "Ho, J., Jain, A., Abbeel, P.",
     "\u201cDenoising Diffusion Probabilistic Models.\u201d \u2014 NeurIPS 2020."),
    (False, "DDIM",
     "Song, J., Meng, C., Ermon, S.",
     "\u201cDenoising Diffusion Implicit Models.\u201d \u2014 ICLR 2021."),
    (True,  "Dataset", "", ""),
    (False, "Mendeley LBC Cervical Cancer",
     "Immagini di citologia cervicale (Pap test), 4 classi diagnostiche.",
     "DOI: 10.17632/zddtpgzv63.2"),
    (True,  "Strumenti Software", "", ""),
]

first = True
for item in entries:
    if item[0]:  # section header
        if first:
            set_paragraph(tf, item[1], size=15, bold=True, space_after=Pt(3), space_before=Pt(0))
            first = False
        else:
            add_paragraph(tf, item[1], size=15, bold=True, space_after=Pt(3), space_before=Pt(10))
    else:  # reference entry
        add_paragraph(tf, item[1], size=13, bold=True, space_after=Pt(0), space_before=Pt(1))
        add_paragraph(tf, f"{item[2]}", size=11, color=COLOR_SUBTLE, space_after=Pt(0))
        add_paragraph(tf, f"{item[3]}", size=11, color=COLOR_SUBTLE, space_after=Pt(4))

# Tools inline
add_paragraph(tf,
    "PyTorch (deep learning) \u2022 scikit-image (PSNR/SSIM) \u2022 Matplotlib (visualizzazioni) \u2022 NumPy/Pandas (dati)",
    size=12, color=COLOR_SUBTLE, space_after=Pt(0), space_before=Pt(2))


# ── Save ──────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Saved to: {OUTPUT}")
