import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
import os

path = r'C:\Users\franc\OneDrive - mapsengineering.com\Università\PrimoAnno\CI\ci-cervical-lbc\slides\presentazione.pptx'
prs = Presentation(path)

print(f'=== FILE: {os.path.basename(path)} ===')
print(f'Slide dimensions: {prs.slide_width}, {prs.slide_height}')
print(f'Total slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides):
    print(f'--- Slide {i+1} ---')
    layout = slide.slide_layout
    print(f'  Layout: {layout.name}')
    
    if slide.placeholders:
        print(f'  Placeholders:')
        for ph in slide.placeholders:
            print(f'    - idx={ph.placeholder_format.idx}, name="{ph.name}", type={ph.placeholder_format.type}')
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'  [{shape.shape_type}] "{text}"')
        if shape.has_table:
            table = shape.table
            print(f'  [TABLE] rows={len(table.rows)}, cols={len(table.columns)}')
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                print(f'    Row {row_idx}: {cells}')
    print()
