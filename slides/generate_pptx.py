"""
Genera la presentazione PowerPoint per l'esame orale di Computational Imaging.
Gruppo R: Francesco Castaldi, Paolo Fusco.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colori ─────────────────────────────────────────────────────────────────────
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2 = RGBColor(0xE0, 0x1E, 0x79)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, y=Inches(0.3), size=Pt(36), color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = True
    return txBox


def add_body(slide, text, y=Inches(1.5), size=Pt(22), color=LIGHT_GRAY, left=Inches(0.8)):
    txBox = slide.shapes.add_textbox(left, y, Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


def add_footer(slide, text="Computational Imaging — Gruppo R | Francesco Castaldi, Paolo Fusco"):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT
    p.font.italic = True


RESULTS_DIR = Path("results")
COMPARISON_PNG = RESULTS_DIR / "comparison.png"


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_title(slide, "Computational Imaging", y=Inches(1.0), size=Pt(48), color=ACCENT)
add_title(slide, "Deblur & Denoise su immagini LBC Cervical", y=Inches(1.8), size=Pt(32), color=WHITE)

add_body(slide, """\u2003Universit\u00e0 di Bologna — LM Informatica
\u2003Prof. Picciolomini & Evangelista
\u2003A.A. 2025/2026
\u2003

\u2003Metodi: Total Variation • UNet • DiffPIR
\u2003Dataset: Mendeley LBC Cervical Cancer (962 immagini)
\u2003

\u2003Fran\u00e7ois Castaldi & Paolo Fusco — Gruppo R""",
          y=Inches(3.0), size=Pt(22))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Task Definition
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Task: Deblur & Denoise", color=ACCENT)

add_body(slide, """Problema inverso mal posto:
    y = (k * x) + n

Degradazione:
  \u2022 Blur Gaussiano: \u03c3 = 2, kernel 9\u00d79
  \u2022 Rumore AWGN: \u03c3\u2099 \u2208 {0.005, 0.01, 0.05, 0.1}
  \u2022 Stessa pipeline degradazione per tutti i metodi (seed=42)

Dataset:
  \u2022 Mendeley LBC Cervical Cancer — 962 immagini, 4 classi diagnostiche
  \u2022 Resize 256\u00d7256, normalizzazione [-1, 1]
  \u2022 Split: 70% train (673) / 15% val (144) / 15% test (145)""",
          y=Inches(1.6), size=Pt(20))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Methods Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "3 Metodi, 3 Famiglie Metodologiche", color=ACCENT)

add_body(slide, """\u2022 Total Variation (TV) — VARIAZIONALE
     \u2514\u2500 minimizza ||H*x - y||\u00b2 + \u03bb\u00b7TV(x) | Adam 150 iter, \u03bb=0.005
     \u2514\u2500 Pro: nessun training, interpretabile, 32 dB a basso rumore
     \u2514\u2500 Contro: \u03bb fisso, staircasing ad alto rumore

\u2022 UNet — END-TO-END DEEP LEARNING
     \u2514\u2500 architettura encoder-decoder, 4 livelli, 31M parametri
     \u2514\u2500 Training: MSE + Adam (lr=1e-4), multi-noise augmentation
     \u2514\u2500 Pro: inferenza rapidissima (26 ms/img), 24 dB con 1 sola epoca
     \u2514\u2500 Contro: richiede training, CPU limita le epoche

\u2022 DiffPIR — GENERATIVO (DIFFUSION)
     \u2514\u2500 DDPM + LightUNet (1.26M param) + FFT data-fidelity
     \u2514\u2500 DDIM sampling, 15 steps, \u03bb=10, t_start=50
     \u2514\u2500 Pro: eccelle ad alto rumore (24.7 dB), ricostruisce dettagli persi
     \u2514\u2500 Contro: lento (2s/img), allucinazioni a basso rumore""",
          y=Inches(1.6), size=Pt(18))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TV Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Risultati TV — Total Variation", color=ACCENT)

# Table
rows, cols = 5, 4
table_shape = slide.shapes.add_table(rows, cols,
    Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.0))
table = table_shape.table

headers = ["Noise \u03c3\u2099", "PSNR (dB)", "SSIM", "Iterazioni"]
data = [
    ["0.005", "32.09", "0.911", "150"],
    ["0.01",  "32.04", "0.909", "150"],
    ["0.05",  "30.42", "0.837", "150"],
    ["0.1",   "26.54", "0.586", "150"],
]

for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = table.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)

add_body(slide, """\u03bb_reg = 0.005 (scelto euristicamente: {0.001, 0.005, 0.01, 0.05, 0.1})
Ottimizzazione Adam, 150 iterazioni, learning rate 0.001

\u2714 Eccelle a basso rumore (32 dB)
\u2716 Cala ad alto rumore (26.5 dB) — \u03bb fisso, staircasing""",
          y=Inches(5.0), size=Pt(16), left=Inches(0.8))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — UNet Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Risultati UNet — End-to-End", color=ACCENT)

table_shape = slide.shapes.add_table(5, 5,
    Inches(0.8), Inches(1.8), Inches(6.5), Inches(3.0))
table = table_shape.table

headers = ["Noise \u03c3\u2099", "PSNR (dB)", "SSIM", "Tempo/img", "Epoche"]
data = [
    ["0.005", "24.07", "0.789", "3.87 s", "1 (CPU)"],
    ["0.01",  "24.05", "0.785", "3.89 s", "1 (CPU)"],
    ["0.05",  "23.45", "0.700", "3.83 s", "1 (CPU)"],
    ["0.1",   "21.87", "0.554", "3.78 s", "1 (CPU)"],
]

for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = table.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)

add_body(slide, """Architettura: encoder-decoder 4 livelli (64\u2192512), skip connections, 31M parametri
Training: multi-noise augmentation (\u03c3\u2099 random per batch), Adam lr=1e-4, batch=16

\u2714 24 dB con 1 sola epoca — potenziale significativo
\u2714 Inferenza rapidissima (~26 ms/img in batch)
\u2716 Limitato a 1 epoca su CPU (~80 min/epoca) — su GPU con 50+ epoche supererebbe TV""",
          y=Inches(5.0), size=Pt(16), left=Inches(0.8))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DiffPIR Results
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Risultati DiffPIR — Generativo (Diffusion)", color=ACCENT)

table_shape = slide.shapes.add_table(5, 5,
    Inches(0.8), Inches(1.8), Inches(6.5), Inches(3.0))
table = table_shape.table

headers = ["Noise \u03c3\u2099", "PSNR (dB)", "SSIM", "Tempo/img", "DDIM steps"]
data = [
    ["0.005", "16.67", "0.235", "2.01 s", "15"],
    ["0.01",  "17.32", "0.270", "1.97 s", "15"],
    ["0.05",  "22.49", "0.512", "2.04 s", "15"],
    ["0.1",   "24.68", "0.664", "2.00 s", "15"],
]

for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = table.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)

add_body(slide, """LightUNet custom (1.26M param) addestrata su LBC. FFT-based data-fidelity.
DDIM sampling: t_start=50, \u03bb=10, \u03b6=0 (deterministico)

\u2714 Eccelle ad alto rumore: 24.7 dB a \u03c3\u2099=0.1 (migliore dei 3 metodi!)
\u2716 PESSIMO a basso rumore: 16.7 dB — il prior generativo allucina dettagli inesistenti
\u2716 Lento: 2 sec/img vs 0.3s (TV) e 0.026s (UNet)""",
          y=Inches(5.0), size=Pt(16), left=Inches(0.8))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Comparison Plot
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Confronto: PSNR & SSIM vs Noise Level", color=ACCENT)

# Insert comparison plot
if COMPARISON_PNG.exists():
    slide.shapes.add_picture(str(COMPARISON_PNG),
        Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0))
else:
    add_body(slide, "(Eseguire python scripts/plot_results.py per generare il grafico)",
              y=Inches(2.5), size=Pt(20))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Tabella Comparativa Completa", color=ACCENT)

table_shape = slide.shapes.add_table(10, 5,
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2))
table = table_shape.table

headers = ["Metodo", "\u03c3\u2099=0.005", "\u03c3\u2099=0.01", "\u03c3\u2099=0.05", "\u03c3\u2099=0.1"]
tv_psnr  = ["TV (PSNR)", "32.09 dB [BEST]", "32.04 dB [BEST]", "30.42 dB [BEST]", "26.54 dB [BEST]"]
tv_ssim  = ["TV (SSIM)", "0.911 [BEST]",  "0.909 [BEST]",  "0.837 [BEST]",  "0.586"]
unet_psnr = ["UNet (PSNR)", "24.07 dB",        "24.05 dB",        "23.45 dB",        "21.87 dB"]
unet_ssim = ["UNet (SSIM)", "0.789",           "0.785",           "0.700",           "0.554"]
diffpir_psnr = ["DiffPIR (PSNR)", "16.67 dB",  "17.32 dB",        "22.49 dB",        "24.68 dB [BEST]"]
diffpir_ssim = ["DiffPIR (SSIM)", "0.235",     "0.270",           "0.512",           "0.664 [BEST]"]

all_data = [tv_psnr, tv_ssim, unet_psnr, unet_ssim, diffpir_psnr, diffpir_ssim]

# Speed row
speed_row = ["\u23f1 Tempo/img", "TV: 0.3 s", "TV: 0.3 s", "TV: 0.3 s", "TV: 0.3 s"]
all_data.append(speed_row)
speed_row2 = ["", "UNet: 0.026 s", "UNet: 0.026 s", "UNet: 0.026 s", "UNet: 0.026 s"]
all_data.append(speed_row2)
speed_row3 = ["", "DiffPIR: 2.0 s", "DiffPIR: 2.0 s", "DiffPIR: 2.0 s", "DiffPIR: 2.0 s"]
all_data.append(speed_row3)

for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

for r, row_data in enumerate(all_data):
    for c, val in enumerate(row_data):
        cell = table.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Conclusions
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Conclusioni: Nessun Metodo Domina su Tutti i Regimi", color=ACCENT)

add_body(slide, """\u2714 BASSO RUMORE (\u03c3\u2099 \u2264 0.01): TV \u00e8 la scelta migliore
      PSNR > 32 dB, nessun training, immediatamente pronto, interpretabile

\u2714 ALTO RUMORE (\u03c3\u2099 = 0.1): DiffPIR \u00e8 il migliore
      24.7 dB PSNR, ricostruisce dettagli che metodi classici non possono recuperare

\u2714 THROUGHPUT: UNet \u00e8 imbattibile in velocit\u00e0
      26 ms/img (100\u00d7 pi\u00f9 veloce di DiffPIR), 24 dB con 1 sola epoca CPU
      Con 50 epoche su GPU ci si aspetta che eguagli o superi TV

\u26a0 LEZIONI APPRESE:
  \u2022 I modelli generativi allucinano a basso rumore (DiffPIR: 16.7 dB peggio dell'input!)
  \u2022 La riproducibilit\u00e0 (stesso seed, stessa degradazione) \u00e8 essenziale per il confronto
  \u2022 Training su CPU \u00e8 il collo di bottiglia principale (1 epoca UNet = 80 min)""",
          y=Inches(1.6), size=Pt(18))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Reproducibility & Code
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_footer(slide)

add_title(slide, "Riproducibilit\u00e0 & Organizzazione", color=ACCENT)

add_body(slide, """\u2714 Stessi input degradati per tutti i metodi (src/degradation/degradation.py)
\u2714 Seed fisso 42 (configs/experiment.yaml)
\u2714 26 test unitari (pytest): degradation, metrics, DiffPIR, UNet
\u2714 4 notebook Jupyter: EDA, TV, UNet, DiffPIR
\u2714 Codice modulare: ogni metodo in src/methods/<metodo>/
\u2714 GitHub: https://github.com/FrancescoCastaldi/ci-cervical-lbc

Parametri scelti euristicamente:
  \u2022 TV: \u03bb_reg=0.005 (testato {0.001, 0.005, 0.01, 0.05, 0.1})
  \u2022 UNet: lr=1e-4, batch=16, multi-noise augmentation
  \u2022 DiffPIR: t_start=50, \u03bb=10, num_steps=15, \u03b6=0

Quick Start:
  pip install -r requirements.txt
  python scripts/run_tv.py
  python scripts/run_unet.py
  python scripts/run_diffpir.py
  python scripts/plot_results.py""",
          y=Inches(1.6), size=Pt(18))


# ── Save ──────────────────────────────────────────────────────────────────────
output_path = Path("slides") / "presentazione.pptx"
prs.save(str(output_path))
print(f"Presentazione salvata: {output_path}")
print(f"Slides: {len(prs.slides)}")
